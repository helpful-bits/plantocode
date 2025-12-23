#!/usr/bin/env python3
"""
Add professional transitions to Vibe Manager presentation
"""

from pptx import Presentation
from pptx.util import Pt
from lxml import etree

def add_transitions_to_presentation(input_path, output_path):
    """
    Add smooth fade transitions to all slides
    PowerPoint transitions are stored in the slide's XML
    """
    prs = Presentation(input_path)

    print(f"\n{'='*60}")
    print(f"Adding Transitions to Presentation")
    print(f"{'='*60}\n")
    print(f"Input: {input_path}")
    print(f"Total slides: {len(prs.slides)}")

    # Transition timing in milliseconds
    TRANSITION_DURATION = 400  # 400ms = 0.4 seconds (smooth and professional)

    for i, slide in enumerate(prs.slides, 1):
        # Access the slide's XML element
        slide_elem = slide.element

        # Create transition element if it doesn't exist
        # In PowerPoint XML, transitions are part of the slide properties
        trans = slide_elem.find('.//{http://schemas.openxmlformats.org/presentationml/2006/main}transition')

        if trans is None:
            # Create new transition element
            trans = etree.Element('{http://schemas.openxmlformats.org/presentationml/2006/main}transition')
            trans.set('spd', 'med')  # Speed: slow, med, fast
            trans.set('dur', str(TRANSITION_DURATION))  # Duration in milliseconds

            # Add fade transition
            fade = etree.SubElement(trans, '{http://schemas.openxmlformats.org/presentationml/2006/main}fade')
            fade.set('thruBlk', '0')  # Fade through black: 0 = no, 1 = yes

            # Insert transition as first child of slide
            slide_elem.insert(0, trans)
        else:
            # Update existing transition
            trans.set('spd', 'med')
            trans.set('dur', str(TRANSITION_DURATION))

            # Remove old transition effects
            for child in list(trans):
                trans.remove(child)

            # Add fade transition
            fade = etree.SubElement(trans, '{http://schemas.openxmlformats.org/presentationml/2006/main}fade')
            fade.set('thruBlk', '0')

        print(f"  [{i:2d}/15] Added fade transition ({TRANSITION_DURATION}ms)")

    # Save modified presentation
    prs.save(output_path)

    print(f"\n{'='*60}")
    print(f"✅ Transitions added successfully!")
    print(f"Output: {output_path}")
    print(f"{'='*60}\n")

    return output_path

if __name__ == '__main__':
    input_file = 'vibe-manager-professional.pptx'
    output_file = 'vibe-manager-professional-with-transitions.pptx'

    add_transitions_to_presentation(input_file, output_file)
