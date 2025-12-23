#!/usr/bin/env python3
"""
Professional Vibe Manager Workflow Presentation
Matches exact design system from desktop app
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# Import exact colors from colors_exact.py
from colors_exact import COLORS

# Assign colors for easy access
BG_GRADIENT_START = COLORS['BG_GRADIENT_START']
BG_NAVY = COLORS['BG_NAVY']
BG_GRADIENT_END = COLORS['BG_GRADIENT_END']
BG_CARD = COLORS['BG_CARD']
BG_ELEVATED = COLORS['BG_ELEVATED']

TEAL_PRIMARY = COLORS['TEAL_PRIMARY']
TEAL_DARK = COLORS['TEAL_DARK']
TEAL_FOREGROUND = COLORS['TEAL_FOREGROUND']

TEXT_PRIMARY = COLORS['TEXT_PRIMARY']
TEXT_SECONDARY = COLORS['TEXT_SECONDARY']
TEXT_MUTED = COLORS['TEXT_MUTED']

BORDER_COLOR = COLORS['BORDER']

# Success/Info colors
SUCCESS_BG = COLORS['SUCCESS_BG']
SUCCESS_COLOR = COLORS['SUCCESS']
INFO_BG = COLORS['INFO_BG']
INFO_COLOR = COLORS['INFO']

def setup_slide_background(slide, gradient=True):
    """Set dark navy background with gradient matching desktop app"""
    background = slide.background
    fill = background.fill
    if gradient:
        fill.gradient()
        fill.gradient_angle = 135.0
        # PowerPoint gradient: 2-stop approximation of 3-stop gradient
        # Desktop uses: linear-gradient(135deg, oklch(0.15 0.02 206) 0%, oklch(0.18 0.02 206) 40%, oklch(0.20 0.03 195) 100%)
        # Approximation: use start and end colors for smooth gradient effect
        fill.gradient_stops[0].color.rgb = BG_GRADIENT_START  # RGB(32, 38, 45) - Darker navy
        fill.gradient_stops[0].position = 0.0
        fill.gradient_stops[1].color.rgb = BG_GRADIENT_END  # RGB(47, 53, 63) - Lighter navy
        fill.gradient_stops[1].position = 1.0
    else:
        fill.solid()
        fill.fore_color.rgb = BG_NAVY

def add_card(slide, left, top, width, height, border_color=None, bg_color=None, glow=False):
    """Add a card shape with proper Vibe Manager styling"""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        left, top, width, height
    )
    shape.adjustments[0] = 0.08  # Corner radius

    # Background
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color or BG_CARD

    # Border
    shape.line.width = Pt(1)
    shape.line.color.rgb = border_color or BORDER_COLOR

    # Subtle shadow for depth
    shadow = shape.shadow
    shadow.inherit = False
    shadow.visible = True
    shadow.shadow_type = 2  # Outer
    shadow.distance = Pt(2)
    shadow.blur_radius = Pt(8)
    shadow.angle = 90
    shadow.transparency = 0.95  # Very subtle

    return shape

def add_title_slide(prs):
    """Slide 1: Hero Title with gradient effect"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    setup_slide_background(slide, gradient=True)

    # Subtitle at top
    subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(0.5))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "INTELLIGENT CODE SCOPE ISOLATION"
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.font.size = Pt(16)
    subtitle_para.font.color.rgb = TEAL_PRIMARY
    subtitle_para.font.bold = True
    subtitle_para.alignment = PP_ALIGN.CENTER
    subtitle_para.font.name = "SF Pro Display"

    # Main title - large and bold
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.text = "Vibe Manager"
    title_frame.word_wrap = True
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(80)
    title_para.font.bold = True
    title_para.font.color.rgb = TEXT_PRIMARY
    title_para.alignment = PP_ALIGN.CENTER
    title_para.font.name = "SF Pro Display"

    # Metric card
    metric_card = add_card(slide, Inches(2.5), Inches(4.5), Inches(5), Inches(1.2))
    metric_text = metric_card.text_frame
    metric_text.word_wrap = True
    metric_text.vertical_anchor = MSO_ANCHOR.MIDDLE

    p1 = metric_text.paragraphs[0]
    p1.text = "1,500 files → 50 files → Production code"
    p1.font.size = Pt(24)
    p1.font.color.rgb = TEXT_SECONDARY
    p1.alignment = PP_ALIGN.CENTER
    p1.font.name = "SF Pro Display"

    p2 = metric_text.add_paragraph()
    p2.text = "In 5 minutes"
    p2.font.size = Pt(32)
    p2.font.bold = True
    p2.font.color.rgb = TEAL_PRIMARY
    p2.alignment = PP_ALIGN.CENTER
    p2.space_before = Pt(8)
    p2.font.name = "SF Pro Display"

def add_section_title(slide, title_text, subtitle_text=None):
    """Add section title with professional styling"""
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.7))
    title_frame = title_box.text_frame
    title_frame.text = title_text
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(40)
    title_para.font.bold = True
    title_para.font.color.rgb = TEXT_PRIMARY
    title_para.font.name = "SF Pro Display"

    if subtitle_text:
        subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.1), Inches(9), Inches(0.4))
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.text = subtitle_text
        subtitle_para = subtitle_frame.paragraphs[0]
        subtitle_para.font.size = Pt(18)
        subtitle_para.font.color.rgb = TEXT_MUTED
        subtitle_para.font.name = "SF Pro Display"
        return 1.7
    return 1.2

def add_problem_slide(prs):
    """Slide 2: The Challenge - Clean layout"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    setup_slide_background(slide)

    start_y = add_section_title(slide, "The Developer's Challenge",
                                 "Working with large-scale codebases")

    # Problem cards in 2x2 grid
    problems = [
        ("Large Codebases", "Navigating 1,500+ files\nwithout clear direction"),
        ("Time Intensive", "Hours spent finding\nrelevant code manually"),
        ("Context Overload", "AI assistants overwhelmed\nby too much information"),
        ("Error Prone", "Manual selection leads to\nmissing critical files")
    ]

    positions = [
        (0.5, 2.2), (5.25, 2.2),
        (0.5, 4.5), (5.25, 4.5)
    ]

    for (title, desc), (x, y) in zip(problems, positions):
        card = add_card(slide, Inches(x), Inches(y), Inches(4.5), Inches(1.8))
        text_frame = card.text_frame
        text_frame.word_wrap = True
        text_frame.margin_left = Inches(0.3)
        text_frame.margin_top = Inches(0.2)

        # Title
        p1 = text_frame.paragraphs[0]
        p1.text = title
        p1.font.size = Pt(22)
        p1.font.bold = True
        p1.font.color.rgb = TEAL_PRIMARY
        p1.font.name = "SF Pro Display"

        # Description
        p2 = text_frame.add_paragraph()
        p2.text = desc
        p2.font.size = Pt(16)
        p2.font.color.rgb = TEXT_SECONDARY
        p2.space_before = Pt(10)
        p2.line_spacing = 1.3
        p2.font.name = "SF Pro Display"

def add_solution_slide(prs):
    """Slide 3: The Solution - Visual flow"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    setup_slide_background(slide)

    add_section_title(slide, "Intelligent Scope Isolation",
                     "From thousands to only what matters")

    # Before state
    before_card = add_card(slide, Inches(1), Inches(2.5), Inches(2.5), Inches(2.5),
                          border_color=TEXT_MUTED)
    before_text = before_card.text_frame
    before_text.vertical_anchor = MSO_ANCHOR.MIDDLE

    p1 = before_text.paragraphs[0]
    p1.text = "1,500+"
    p1.font.size = Pt(56)
    p1.font.bold = True
    p1.font.color.rgb = TEXT_MUTED
    p1.alignment = PP_ALIGN.CENTER
    p1.font.name = "SF Pro Display"

    p2 = before_text.add_paragraph()
    p2.text = "files in codebase"
    p2.font.size = Pt(18)
    p2.font.color.rgb = TEXT_MUTED
    p2.alignment = PP_ALIGN.CENTER
    p2.space_before = Pt(12)
    p2.font.name = "SF Pro Display"

    # Arrow with process label
    arrow_box = slide.shapes.add_textbox(Inches(3.8), Inches(3.5), Inches(2.4), Inches(0.6))
    arrow_text = arrow_box.text_frame
    arrow_text.vertical_anchor = MSO_ANCHOR.MIDDLE

    arrow_para = arrow_text.paragraphs[0]
    arrow_para.text = "→"
    arrow_para.font.size = Pt(72)
    arrow_para.font.color.rgb = TEAL_PRIMARY
    arrow_para.alignment = PP_ALIGN.CENTER
    arrow_para.font.name = "SF Pro Display"

    # Process label
    process_box = slide.shapes.add_textbox(Inches(3.8), Inches(4.2), Inches(2.4), Inches(0.4))
    process_text = process_box.text_frame
    process_para = process_text.paragraphs[0]
    process_para.text = "AI-Powered Analysis"
    process_para.font.size = Pt(14)
    process_para.font.color.rgb = TEAL_PRIMARY
    process_para.alignment = PP_ALIGN.CENTER
    process_para.font.name = "SF Pro Display"

    # After state
    after_card = add_card(slide, Inches(6.5), Inches(2.5), Inches(2.5), Inches(2.5),
                         border_color=TEAL_PRIMARY)
    after_text = after_card.text_frame
    after_text.vertical_anchor = MSO_ANCHOR.MIDDLE

    p3 = after_text.paragraphs[0]
    p3.text = "50"
    p3.font.size = Pt(56)
    p3.font.bold = True
    p3.font.color.rgb = TEAL_PRIMARY
    p3.alignment = PP_ALIGN.CENTER
    p3.font.name = "SF Pro Display"

    p4 = after_text.add_paragraph()
    p4.text = "relevant files"
    p4.font.size = Pt(18)
    p4.font.color.rgb = TEXT_SECONDARY
    p4.alignment = PP_ALIGN.CENTER
    p4.space_before = Pt(12)
    p4.font.name = "SF Pro Display"

    # Result metric at bottom
    result_card = add_card(slide, Inches(2.5), Inches(5.5), Inches(5), Inches(0.8))
    result_text = result_card.text_frame
    result_text.vertical_anchor = MSO_ANCHOR.MIDDLE

    result_para = result_text.paragraphs[0]
    result_para.text = "96.7% reduction in scope | 100% accuracy"
    result_para.font.size = Pt(20)
    result_para.font.color.rgb = SUCCESS_COLOR
    result_para.alignment = PP_ALIGN.CENTER
    result_para.font.name = "SF Pro Display"

def add_parallel_workflow_slide(prs):
    """Slide 4: Parallel File Discovery"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    setup_slide_background(slide)

    add_section_title(slide, "Parallel File Discovery",
                     "Multiple AI agents searching simultaneously")

    # Source
    source_card = add_card(slide, Inches(0.8), Inches(2.2), Inches(2), Inches(1.2))
    source_text = source_card.text_frame
    source_text.vertical_anchor = MSO_ANCHOR.MIDDLE

    source_para = source_text.paragraphs[0]
    source_para.text = "Codebase\nAnalysis"
    source_para.font.size = Pt(20)
    source_para.font.bold = True
    source_para.font.color.rgb = TEXT_PRIMARY
    source_para.alignment = PP_ALIGN.CENTER
    source_para.line_spacing = 1.2
    source_para.font.name = "SF Pro Display"

    # Parallel workflows
    workflows = [
        "Semantic Search",
        "Pattern Matching",
        "Dependency Analysis",
        "AST Traversal"
    ]

    y_positions = [2, 2.9, 3.8, 4.7]

    for workflow, y in zip(workflows, y_positions):
        # Arrow
        arrow_shape = slide.shapes.add_shape(
            MSO_SHAPE.RIGHT_ARROW,
            Inches(3), Inches(y), Inches(1.8), Inches(0.5)
        )
        arrow_shape.fill.solid()
        arrow_shape.fill.fore_color.rgb = TEAL_DARK
        arrow_shape.line.color.rgb = TEAL_PRIMARY
        arrow_shape.line.width = Pt(1)

        # Workflow card
        workflow_card = add_card(slide, Inches(5), Inches(y), Inches(3.5), Inches(0.7))
        workflow_text = workflow_card.text_frame
        workflow_text.vertical_anchor = MSO_ANCHOR.MIDDLE
        workflow_text.margin_left = Inches(0.2)

        workflow_para = workflow_text.paragraphs[0]
        workflow_para.text = workflow
        workflow_para.font.size = Pt(18)
        workflow_para.font.color.rgb = TEXT_PRIMARY
        workflow_para.font.name = "SF Pro Display"

    # Result
    result_card = add_card(slide, Inches(3), Inches(5.7), Inches(5.5), Inches(0.8),
                          border_color=TEAL_PRIMARY)
    result_text = result_card.text_frame
    result_text.vertical_anchor = MSO_ANCHOR.MIDDLE

    result_para = result_text.paragraphs[0]
    result_para.text = "Converged Results: 50 High-Confidence Files"
    result_para.font.size = Pt(22)
    result_para.font.bold = True
    result_para.font.color.rgb = TEAL_PRIMARY
    result_para.alignment = PP_ALIGN.CENTER
    result_para.font.name = "SF Pro Display"

def add_metrics_slide(prs):
    """Slide 5: Key Metrics - Dashboard style"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    setup_slide_background(slide)

    add_section_title(slide, "Precision Metrics",
                     "Quantifiable scope reduction and accuracy")

    # Main metrics - 2x2 grid
    metrics = [
        ("1,500", "Original\nFiles", TEXT_MUTED),
        ("50", "Discovered\nFiles", TEAL_PRIMARY),
        ("3.3%", "Final\nScope", SUCCESS_COLOR),
        ("96%+", "Accuracy\nRate", INFO_COLOR)
    ]

    positions = [
        (0.8, 2.2), (5.2, 2.2),
        (0.8, 4.2), (5.2, 4.2)
    ]

    for (number, label, color), (x, y) in zip(metrics, positions):
        card = add_card(slide, Inches(x), Inches(y), Inches(4), Inches(1.6))
        text_frame = card.text_frame
        text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

        # Number
        p1 = text_frame.paragraphs[0]
        p1.text = number
        p1.font.size = Pt(52)
        p1.font.bold = True
        p1.font.color.rgb = color
        p1.alignment = PP_ALIGN.CENTER
        p1.font.name = "SF Pro Display"

        # Label
        p2 = text_frame.add_paragraph()
        p2.text = label
        p2.font.size = Pt(18)
        p2.font.color.rgb = TEXT_MUTED
        p2.alignment = PP_ALIGN.CENTER
        p2.space_before = Pt(8)
        p2.line_spacing = 1.2
        p2.font.name = "SF Pro Display"

def add_verification_slide(prs):
    """Slide 6: User Verification"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    setup_slide_background(slide)

    add_section_title(slide, "Human-in-the-Loop Verification",
                     "Maintain control with manual review")

    # Main content card
    content_card = add_card(slide, Inches(1.5), Inches(2.2), Inches(7), Inches(3.6))
    text_frame = content_card.text_frame
    text_frame.word_wrap = True
    text_frame.margin_left = Inches(0.4)
    text_frame.margin_top = Inches(0.3)

    # Title
    p_title = text_frame.paragraphs[0]
    p_title.text = "Review and Adjust Discovered Files"
    p_title.font.size = Pt(26)
    p_title.font.bold = True
    p_title.font.color.rgb = TEXT_PRIMARY
    p_title.font.name = "SF Pro Display"

    # Capabilities
    capabilities = [
        "Remove irrelevant or redundant files",
        "Add additional context files manually",
        "Verify scope completeness before proceeding",
        "Adjust search parameters if needed"
    ]

    for capability in capabilities:
        p = text_frame.add_paragraph()
        p.text = f"•  {capability}"
        p.font.size = Pt(20)
        p.font.color.rgb = TEXT_SECONDARY
        p.space_before = Pt(16)
        p.line_spacing = 1.3
        p.font.name = "SF Pro Display"

    # Quality assurance note
    qa_card = add_card(slide, Inches(2.5), Inches(6), Inches(5), Inches(0.7),
                      border_color=INFO_COLOR, bg_color=INFO_BG)
    qa_text = qa_card.text_frame
    qa_text.vertical_anchor = MSO_ANCHOR.MIDDLE

    qa_para = qa_text.paragraphs[0]
    qa_para.text = "Ensures precision while maintaining automation benefits"
    qa_para.font.size = Pt(18)
    qa_para.font.italic = True
    qa_para.font.color.rgb = INFO_COLOR
    qa_para.alignment = PP_ALIGN.CENTER
    qa_para.font.name = "SF Pro Display"

def add_xml_slide(prs):
    """Slide 7: XML Structured Format"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    setup_slide_background(slide)

    add_section_title(slide, "XML-Structured Communication",
                     "Optimized format for LLM processing")

    # Benefits
    benefits_card = add_card(slide, Inches(0.8), Inches(2), Inches(4), Inches(2))
    benefits_text = benefits_card.text_frame
    benefits_text.word_wrap = True
    benefits_text.margin_left = Inches(0.3)
    benefits_text.margin_top = Inches(0.2)

    p_title = benefits_text.paragraphs[0]
    p_title.text = "Why XML?"
    p_title.font.size = Pt(24)
    p_title.font.bold = True
    p_title.font.color.rgb = TEAL_PRIMARY
    p_title.font.name = "SF Pro Display"

    benefits = [
        "Structured & parseable",
        "LLM-optimized format",
        "Hierarchical context",
        "Version traceable"
    ]

    for benefit in benefits:
        p = benefits_text.add_paragraph()
        p.text = f"•  {benefit}"
        p.font.size = Pt(18)
        p.font.color.rgb = TEXT_SECONDARY
        p.space_before = Pt(12)
        p.font.name = "SF Pro Display"

    # Code example
    code_card = add_card(slide, Inches(5.2), Inches(2), Inches(4), Inches(2),
                        bg_color=BG_ELEVATED)
    code_text = code_card.text_frame
    code_text.word_wrap = True
    code_text.margin_left = Inches(0.25)
    code_text.margin_top = Inches(0.2)

    code_para = code_text.paragraphs[0]
    code_content = """<context>
  <task>Feature X</task>
  <files>
    <file path="src/...">
      <!-- content -->
    </file>
  </files>
</context>"""
    code_para.text = code_content
    code_para.font.name = "Monaco"
    code_para.font.size = Pt(15)
    code_para.font.color.rgb = TEAL_PRIMARY
    code_para.line_spacing = 1.4

    # Model integration
    model_card = add_card(slide, Inches(2), Inches(4.5), Inches(6), Inches(0.9),
                         border_color=TEAL_PRIMARY)
    model_text = model_card.text_frame
    model_text.vertical_anchor = MSO_ANCHOR.MIDDLE

    model_para = model_text.paragraphs[0]
    model_para.text = "Optimized for OpenAI GPT-4 and Claude 3.5 Sonnet"
    model_para.font.size = Pt(20)
    model_para.font.bold = True
    model_para.font.color.rgb = TEXT_PRIMARY
    model_para.alignment = PP_ALIGN.CENTER
    model_para.font.name = "SF Pro Display"

    sub_para = model_text.add_paragraph()
    sub_para.text = "Regularly updated to leverage latest model capabilities"
    sub_para.font.size = Pt(16)
    sub_para.font.color.rgb = TEXT_MUTED
    sub_para.alignment = PP_ALIGN.CENTER
    sub_para.space_before = Pt(6)
    sub_para.font.name = "SF Pro Display"

def add_triple_planning_slide(prs):
    """Slide 8: Triple Parallel Planning"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    setup_slide_background(slide)

    add_section_title(slide, "Triple Parallel Planning",
                     "Three implementation strategies generated concurrently")

    # Input
    input_card = add_card(slide, Inches(0.8), Inches(2.2), Inches(2.2), Inches(1.2))
    input_text = input_card.text_frame
    input_text.vertical_anchor = MSO_ANCHOR.MIDDLE

    input_para = input_text.paragraphs[0]
    input_para.text = "Verified\nScope (50)"
    input_para.font.size = Pt(20)
    input_para.font.bold = True
    input_para.font.color.rgb = TEXT_PRIMARY
    input_para.alignment = PP_ALIGN.CENTER
    input_para.line_spacing = 1.2
    input_para.font.name = "SF Pro Display"

    # Three planning agents
    plans = ["Plan A", "Plan B", "Plan C"]
    y_positions = [2.2, 3.6, 5.0]

    for i, (plan, y) in enumerate(zip(plans, y_positions)):
        # Arrow
        arrow = slide.shapes.add_shape(
            MSO_SHAPE.RIGHT_ARROW,
            Inches(3.2), Inches(y), Inches(1.5), Inches(0.5)
        )
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = TEAL_DARK
        arrow.line.color.rgb = TEAL_PRIMARY
        arrow.line.width = Pt(1)

        # Plan card
        plan_card = add_card(slide, Inches(5), Inches(y), Inches(4), Inches(1.1))
        plan_text = plan_card.text_frame
        plan_text.vertical_anchor = MSO_ANCHOR.MIDDLE
        plan_text.margin_left = Inches(0.3)

        plan_title = plan_text.paragraphs[0]
        plan_title.text = f"Implementation {plan}"
        plan_title.font.size = Pt(22)
        plan_title.font.bold = True
        plan_title.font.color.rgb = TEAL_PRIMARY
        plan_title.font.name = "SF Pro Display"

        plan_desc = plan_text.add_paragraph()
        plan_desc.text = f"Independent strategy with unique approach"
        plan_desc.font.size = Pt(16)
        plan_desc.font.color.rgb = TEXT_MUTED
        plan_desc.space_before = Pt(6)
        plan_desc.font.name = "SF Pro Display"

    # Benefit
    benefit_card = add_card(slide, Inches(2), Inches(6.3), Inches(6), Inches(0.6))
    benefit_text = benefit_card.text_frame
    benefit_text.vertical_anchor = MSO_ANCHOR.MIDDLE

    benefit_para = benefit_text.paragraphs[0]
    benefit_para.text = "Parallel execution reduces planning time by 67%"
    benefit_para.font.size = Pt(18)
    benefit_para.font.color.rgb = SUCCESS_COLOR
    benefit_para.alignment = PP_ALIGN.CENTER
    benefit_para.font.name = "SF Pro Display"

def add_scope_narrowing_slide(prs):
    """Slide 9: Second Scope Narrowing"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    setup_slide_background(slide)

    add_section_title(slide, "Precise Change Identification",
                     "From 50 files to specific modifications")

    # Funnel visualization
    funnel_card = add_card(slide, Inches(2), Inches(2.2), Inches(6), Inches(3.5))
    text_frame = funnel_card.text_frame
    text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

    # Stage 1
    p1 = text_frame.paragraphs[0]
    p1.text = "50 Files Analyzed"
    p1.font.size = Pt(32)
    p1.font.bold = True
    p1.font.color.rgb = TEXT_MUTED
    p1.alignment = PP_ALIGN.CENTER
    p1.font.name = "SF Pro Display"

    # Arrow
    p_arrow1 = text_frame.add_paragraph()
    p_arrow1.text = "↓"
    p_arrow1.font.size = Pt(48)
    p_arrow1.font.color.rgb = TEAL_PRIMARY
    p_arrow1.alignment = PP_ALIGN.CENTER
    p_arrow1.space_before = Pt(8)
    p_arrow1.font.name = "SF Pro Display"

    # Stage 2
    p2 = text_frame.add_paragraph()
    p2.text = "8-10 Files to Modify"
    p2.font.size = Pt(28)
    p2.font.bold = True
    p2.font.color.rgb = TEAL_PRIMARY
    p2.alignment = PP_ALIGN.CENTER
    p2.space_before = Pt(8)
    p2.font.name = "SF Pro Display"

    # Arrow
    p_arrow2 = text_frame.add_paragraph()
    p_arrow2.text = "↓"
    p_arrow2.font.size = Pt(48)
    p_arrow2.font.color.rgb = TEAL_PRIMARY
    p_arrow2.alignment = PP_ALIGN.CENTER
    p_arrow2.space_before = Pt(8)
    p_arrow2.font.name = "SF Pro Display"

    # Stage 3
    p3 = text_frame.add_paragraph()
    p3.text = "4 Files to Create"
    p3.font.size = Pt(24)
    p3.font.bold = True
    p3.font.color.rgb = SUCCESS_COLOR
    p3.alignment = PP_ALIGN.CENTER
    p3.space_before = Pt(8)
    p3.font.name = "SF Pro Display"

    # Precision metric
    precision_card = add_card(slide, Inches(2.5), Inches(6), Inches(5), Inches(0.7))
    precision_text = precision_card.text_frame
    precision_text.vertical_anchor = MSO_ANCHOR.MIDDLE

    precision_para = precision_text.paragraphs[0]
    precision_para.text = "94% Precision in Change Detection"
    precision_para.font.size = Pt(22)
    precision_para.font.bold = True
    precision_para.font.color.rgb = SUCCESS_COLOR
    precision_para.alignment = PP_ALIGN.CENTER
    precision_para.font.name = "SF Pro Display"

def add_judge_architect_slide(prs):
    """Slide 10: Judge Architect"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    setup_slide_background(slide)

    add_section_title(slide, "Judge Architect",
                     "AI-powered plan synthesis and optimization")

    # Input sources
    sources_card = add_card(slide, Inches(0.8), Inches(2.2), Inches(3.8), Inches(3))
    sources_text = sources_card.text_frame
    sources_text.word_wrap = True
    sources_text.margin_left = Inches(0.3)
    sources_text.margin_top = Inches(0.25)

    sources_title = sources_text.paragraphs[0]
    sources_title.text = "Input Sources"
    sources_title.font.size = Pt(24)
    sources_title.font.bold = True
    sources_title.font.color.rgb = TEAL_PRIMARY
    sources_title.font.name = "SF Pro Display"

    sources_list = [
        "Selected implementation plans",
        "Mentioned files from plans",
        "Project directory tree",
        "User extra instructions",
        "Dependency graph"
    ]

    for source in sources_list:
        p = sources_text.add_paragraph()
        p.text = f"•  {source}"
        p.font.size = Pt(17)
        p.font.color.rgb = TEXT_SECONDARY
        p.space_before = Pt(10)
        p.font.name = "SF Pro Display"

    # Process
    process_box = slide.shapes.add_textbox(Inches(4.8), Inches(3.5), Inches(0.6), Inches(0.6))
    process_text = process_box.text_frame
    process_text.vertical_anchor = MSO_ANCHOR.MIDDLE

    process_para = process_text.paragraphs[0]
    process_para.text = "→"
    process_para.font.size = Pt(60)
    process_para.font.color.rgb = TEAL_PRIMARY
    process_para.alignment = PP_ALIGN.CENTER
    process_para.font.name = "SF Pro Display"

    # Output
    output_card = add_card(slide, Inches(5.6), Inches(2.5), Inches(3.6), Inches(2.5),
                          border_color=TEAL_PRIMARY)
    output_text = output_card.text_frame
    output_text.vertical_anchor = MSO_ANCHOR.MIDDLE

    output_title = output_text.paragraphs[0]
    output_title.text = "Unified Plan"
    output_title.font.size = Pt(28)
    output_title.font.bold = True
    output_title.font.color.rgb = TEAL_PRIMARY
    output_title.alignment = PP_ALIGN.CENTER
    output_title.font.name = "SF Pro Display"

    output_desc = output_text.add_paragraph()
    output_desc.text = "Optimized synthesis\nof best approaches"
    output_desc.font.size = Pt(18)
    output_desc.font.color.rgb = TEXT_SECONDARY
    output_desc.alignment = PP_ALIGN.CENTER
    output_desc.space_before = Pt(12)
    output_desc.line_spacing = 1.3
    output_desc.font.name = "SF Pro Display"

    # Quality note
    quality_card = add_card(slide, Inches(1.5), Inches(5.5), Inches(7), Inches(0.9))
    quality_text = quality_card.text_frame
    quality_text.vertical_anchor = MSO_ANCHOR.MIDDLE
    quality_text.margin_left = Inches(0.3)

    quality_para = quality_text.paragraphs[0]
    quality_para.text = "Evaluates consistency, completeness, and architectural soundness"
    quality_para.font.size = Pt(20)
    quality_para.font.color.rgb = TEXT_SECONDARY
    quality_para.alignment = PP_ALIGN.CENTER
    quality_para.font.name = "SF Pro Display"

def add_multi_agent_slide(prs):
    """Slide 11: Multi-Agent Execution"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    setup_slide_background(slide)

    add_section_title(slide, "Parallel Code Generation",
                     "Multiple specialized agents executing simultaneously")

    # Agent cards in grid
    agents = [
        "Authentication\nModule",
        "UI Components\n& Styling",
        "API Endpoints\n& Routes",
        "Database\nMigrations",
        "Unit Tests\n& Integration",
        "Documentation\n& Types"
    ]

    positions = [
        (0.8, 2.2), (3.7, 2.2), (6.6, 2.2),
        (0.8, 4), (3.7, 4), (6.6, 4)
    ]

    for agent, (x, y) in zip(agents, positions):
        card = add_card(slide, Inches(x), Inches(y), Inches(2.6), Inches(1.4))
        text_frame = card.text_frame
        text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

        agent_para = text_frame.paragraphs[0]
        agent_para.text = agent
        agent_para.font.size = Pt(17)
        agent_para.font.bold = True
        agent_para.font.color.rgb = TEXT_PRIMARY
        agent_para.alignment = PP_ALIGN.CENTER
        agent_para.line_spacing = 1.3
        agent_para.font.name = "SF Pro Display"

    # Timeline
    timeline_card = add_card(slide, Inches(2), Inches(5.8), Inches(6), Inches(0.9),
                            border_color=SUCCESS_COLOR, bg_color=SUCCESS_BG)
    timeline_text = timeline_card.text_frame
    timeline_text.vertical_anchor = MSO_ANCHOR.MIDDLE

    timeline_para = timeline_text.paragraphs[0]
    timeline_para.text = "Complete Feature Implementation"
    timeline_para.font.size = Pt(24)
    timeline_para.font.bold = True
    timeline_para.font.color.rgb = SUCCESS_COLOR
    timeline_para.alignment = PP_ALIGN.CENTER
    timeline_para.font.name = "SF Pro Display"

    timeline_sub = timeline_text.add_paragraph()
    timeline_sub.text = "5 minutes average execution time"
    timeline_sub.font.size = Pt(18)
    timeline_sub.font.color.rgb = TEXT_MUTED
    timeline_sub.alignment = PP_ALIGN.CENTER
    timeline_sub.space_before = Pt(6)
    timeline_sub.font.name = "SF Pro Display"

def add_results_slide(prs):
    """Slide 12: Production Results"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    setup_slide_background(slide)

    add_section_title(slide, "Production-Ready Results",
                     "High-quality code in minutes, not days")

    # Results grid
    results = [
        ("90%+", "Completion\nRate", SUCCESS_COLOR),
        ("0", "Regression\nBugs", SUCCESS_COLOR),
        ("< 5%", "Compile\nErrors", INFO_COLOR),
        ("100%", "Type\nSafety", SUCCESS_COLOR)
    ]

    positions = [
        (0.8, 2.2), (5.2, 2.2),
        (0.8, 4.2), (5.2, 4.2)
    ]

    for (metric, label, color), (x, y) in zip(results, positions):
        card = add_card(slide, Inches(x), Inches(y), Inches(4), Inches(1.6))
        text_frame = card.text_frame
        text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

        metric_para = text_frame.paragraphs[0]
        metric_para.text = metric
        metric_para.font.size = Pt(52)
        metric_para.font.bold = True
        metric_para.font.color.rgb = color
        metric_para.alignment = PP_ALIGN.CENTER
        metric_para.font.name = "SF Pro Display"

        label_para = text_frame.add_paragraph()
        label_para.text = label
        label_para.font.size = Pt(18)
        label_para.font.color.rgb = TEXT_MUTED
        label_para.alignment = PP_ALIGN.CENTER
        label_para.space_before = Pt(8)
        label_para.line_spacing = 1.2
        label_para.font.name = "SF Pro Display"

    # Quality note
    quality_card = add_card(slide, Inches(1.5), Inches(6.1), Inches(7), Inches(0.6))
    quality_text = quality_card.text_frame
    quality_text.vertical_anchor = MSO_ANCHOR.MIDDLE

    quality_para = quality_text.paragraphs[0]
    quality_para.text = "Minor polish occasionally needed | Maintains code quality standards"
    quality_para.font.size = Pt(17)
    quality_para.font.color.rgb = TEXT_MUTED
    quality_para.alignment = PP_ALIGN.CENTER
    quality_para.font.name = "SF Pro Display"

def add_comparison_slide(prs):
    """Slide 13: Traditional vs Vibe Manager"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    setup_slide_background(slide)

    add_section_title(slide, "Development Workflow Comparison",
                     "Traditional approach vs. AI-powered automation")

    # Traditional
    trad_card = add_card(slide, Inches(0.6), Inches(2.2), Inches(4.4), Inches(4),
                        border_color=TEXT_MUTED)
    trad_text = trad_card.text_frame
    trad_text.word_wrap = True
    trad_text.margin_left = Inches(0.3)
    trad_text.margin_top = Inches(0.25)

    trad_title = trad_text.paragraphs[0]
    trad_title.text = "Traditional Workflow"
    trad_title.font.size = Pt(24)
    trad_title.font.bold = True
    trad_title.font.color.rgb = TEXT_MUTED
    trad_title.font.name = "SF Pro Display"

    trad_points = [
        "Hours of manual file hunting",
        "Context overload and fatigue",
        "Higher error rates",
        "Single-threaded execution",
        "2-3 days per feature",
        "Incomplete coverage risk"
    ]

    for point in trad_points:
        p = trad_text.add_paragraph()
        p.text = f"•  {point}"
        p.font.size = Pt(17)
        p.font.color.rgb = TEXT_MUTED
        p.space_before = Pt(12)
        p.font.name = "SF Pro Display"

    # Vibe Manager
    vm_card = add_card(slide, Inches(5.4), Inches(2.2), Inches(4.4), Inches(4),
                      border_color=TEAL_PRIMARY)
    vm_text = vm_card.text_frame
    vm_text.word_wrap = True
    vm_text.margin_left = Inches(0.3)
    vm_text.margin_top = Inches(0.25)

    vm_title = vm_text.paragraphs[0]
    vm_title.text = "Vibe Manager"
    vm_title.font.size = Pt(24)
    vm_title.font.bold = True
    vm_title.font.color.rgb = TEAL_PRIMARY
    vm_title.font.name = "SF Pro Display"

    vm_points = [
        "Seconds to identify scope",
        "Precise context isolation",
        "AI-verified accuracy",
        "Parallel multi-agent execution",
        "5 minutes per feature",
        "Complete test coverage"
    ]

    for point in vm_points:
        p = vm_text.add_paragraph()
        p.text = f"•  {point}"
        p.font.size = Pt(17)
        p.font.color.rgb = TEXT_PRIMARY
        p.space_before = Pt(12)
        p.font.name = "SF Pro Display"

def add_differentiators_slide(prs):
    """Slide 14: Key Differentiators"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    setup_slide_background(slide)

    add_section_title(slide, "What Sets Us Apart",
                     "Core technological advantages")

    differentiators = [
        ("Intelligent Scope Isolation",
         "Multi-stage reduction from 1,500 to 12 files with 96%+ accuracy through parallel AI workflows"),
        ("XML-Powered Precision",
         "Structured communication optimized for LLMs with full parseability and traceability"),
        ("Triple Planning Architecture",
         "Three independent implementation strategies merged by Judge Architect for optimal solution"),
        ("Parallel Multi-Agent Execution",
         "Specialized agents working simultaneously across different code domains for 10x speed increase")
    ]

    y_positions = [2, 3.15, 4.3, 5.45]

    for (title, desc), y in zip(differentiators, y_positions):
        card = add_card(slide, Inches(0.8), Inches(y), Inches(8.4), Inches(1))
        text_frame = card.text_frame
        text_frame.word_wrap = True
        text_frame.margin_left = Inches(0.3)
        text_frame.margin_top = Inches(0.15)
        text_frame.margin_right = Inches(0.3)

        title_para = text_frame.paragraphs[0]
        title_para.text = title
        title_para.font.size = Pt(20)
        title_para.font.bold = True
        title_para.font.color.rgb = TEAL_PRIMARY
        title_para.font.name = "SF Pro Display"

        desc_para = text_frame.add_paragraph()
        desc_para.text = desc
        desc_para.font.size = Pt(15)
        desc_para.font.color.rgb = TEXT_SECONDARY
        desc_para.space_before = Pt(6)
        desc_para.font.name = "SF Pro Display"

def add_cta_slide(prs):
    """Slide 15: Call to Action"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    setup_slide_background(slide, gradient=True)

    # Main title
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = "Transform Your Development Workflow"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(44)
    title_para.font.bold = True
    title_para.font.color.rgb = TEXT_PRIMARY
    title_para.alignment = PP_ALIGN.CENTER
    title_para.font.name = "SF Pro Display"

    # Metrics recap
    metrics_card = add_card(slide, Inches(2), Inches(3.3), Inches(6), Inches(1.2),
                           border_color=TEAL_PRIMARY)
    metrics_text = metrics_card.text_frame
    metrics_text.vertical_anchor = MSO_ANCHOR.MIDDLE

    metrics_para = metrics_text.paragraphs[0]
    metrics_para.text = "1,500 files → 50 files → Production code"
    metrics_para.font.size = Pt(28)
    metrics_para.font.bold = True
    metrics_para.font.color.rgb = TEXT_PRIMARY
    metrics_para.alignment = PP_ALIGN.CENTER
    metrics_para.font.name = "SF Pro Display"

    metrics_sub = metrics_text.add_paragraph()
    metrics_sub.text = "In just 5 minutes"
    metrics_sub.font.size = Pt(24)
    metrics_sub.font.color.rgb = TEAL_PRIMARY
    metrics_sub.alignment = PP_ALIGN.CENTER
    metrics_sub.space_before = Pt(8)
    metrics_sub.font.name = "SF Pro Display"

    # Contact info
    contact_card = add_card(slide, Inches(2.5), Inches(5), Inches(5), Inches(1.2))
    contact_text = contact_card.text_frame
    contact_text.vertical_anchor = MSO_ANCHOR.MIDDLE

    website_para = contact_text.paragraphs[0]
    website_para.text = "vibemanager.app"
    website_para.font.size = Pt(32)
    website_para.font.bold = True
    website_para.font.color.rgb = TEAL_PRIMARY
    website_para.alignment = PP_ALIGN.CENTER
    website_para.font.name = "SF Pro Display"

    email_para = contact_text.add_paragraph()
    email_para.text = "hello@vibemanager.app"
    email_para.font.size = Pt(20)
    email_para.font.color.rgb = TEXT_SECONDARY
    email_para.alignment = PP_ALIGN.CENTER
    email_para.space_before = Pt(8)
    email_para.font.name = "SF Pro Display"

def main():
    """Generate the professional presentation"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    print("\nGenerating Professional Vibe Manager Presentation...")
    print("=" * 60)

    slides = [
        ("Hero Title", add_title_slide),
        ("Problem Statement", add_problem_slide),
        ("Solution Overview", add_solution_slide),
        ("Parallel Workflow", add_parallel_workflow_slide),
        ("Key Metrics", add_metrics_slide),
        ("Verification", add_verification_slide),
        ("XML Format", add_xml_slide),
        ("Triple Planning", add_triple_planning_slide),
        ("Scope Narrowing", add_scope_narrowing_slide),
        ("Judge Architect", add_judge_architect_slide),
        ("Multi-Agent Execution", add_multi_agent_slide),
        ("Results", add_results_slide),
        ("Comparison", add_comparison_slide),
        ("Differentiators", add_differentiators_slide),
        ("Call to Action", add_cta_slide),
    ]

    for i, (name, func) in enumerate(slides, 1):
        print(f"  [{i:2d}/15] Creating: {name}")
        func(prs)

    # Save
    output_path = 'vibe-manager-professional.pptx'
    prs.save(output_path)

    print("=" * 60)
    print(f"\n✅ Professional presentation created successfully!")
    print(f"   Location: {output_path}")
    print(f"   Slides: {len(prs.slides)}")
    print(f"   File size: ~{len(prs.slides) * 3}KB (estimated)")
    print("\n" + "=" * 60)

if __name__ == '__main__':
    main()
