#!/usr/bin/env python3
"""
Generate professional marketing video from PowerPoint presentation
Exports slides as high-quality images and creates video with FFmpeg
"""

import os
import subprocess
from pptx import Presentation
from PIL import Image, ImageDraw, ImageFont
import io

# Video configuration
VIDEO_CONFIG = {
    'width': 1920,
    'height': 1080,
    'fps': 30,
    'slide_duration': 8,  # seconds per slide
    'transition_duration': 0.5,  # seconds for transitions
    'output_format': 'mp4',
    'codec': 'libx264',
    'quality': 'high',  # high, medium, low
    'bitrate': '8000k',
}

def check_ffmpeg():
    """Check if FFmpeg is installed"""
    try:
        result = subprocess.run(['ffmpeg', '-version'],
                              capture_output=True,
                              text=True)
        return result.returncode == 0
    except FileNotFoundError:
        return False

def export_slides_as_images(pptx_path, output_dir='slides_export'):
    """
    Export each slide as high-resolution PNG image
    Note: This requires LibreOffice or similar converter
    """
    os.makedirs(output_dir, exist_ok=True)

    print(f"\n{'='*60}")
    print(f"Exporting Slides as Images")
    print(f"{'='*60}\n")

    # Check if LibreOffice is available
    libreoffice_paths = [
        '/Applications/LibreOffice.app/Contents/MacOS/soffice',  # macOS
        '/usr/bin/soffice',  # Linux
        'soffice',  # Windows/PATH
    ]

    soffice_path = None
    for path in libreoffice_paths:
        if os.path.exists(path) or subprocess.run(['which', path.split('/')[-1]],
                                                   capture_output=True).returncode == 0:
            soffice_path = path
            break

    if not soffice_path:
        print("⚠️  LibreOffice not found. Using alternative method...")
        return export_slides_manual(pptx_path, output_dir)

    # Convert PPTX to PDF first (better quality)
    pdf_path = pptx_path.replace('.pptx', '.pdf')
    cmd = [
        soffice_path,
        '--headless',
        '--convert-to', 'pdf',
        '--outdir', os.path.dirname(pptx_path) or '.',
        pptx_path
    ]

    print(f"Converting presentation to PDF...")
    result = subprocess.run(cmd, capture_output=True, text=True)

    if result.returncode != 0 or not os.path.exists(pdf_path):
        print("⚠️  PDF conversion failed. Using alternative method...")
        return export_slides_manual(pptx_path, output_dir)

    # Convert PDF to images using ImageMagick or similar
    print(f"Converting PDF pages to images...")

    # Try using ImageMagick's convert command
    try:
        cmd = [
            'convert',
            '-density', '300',  # High DPI
            '-quality', '100',
            pdf_path,
            f'{output_dir}/slide_%02d.png'
        ]
        subprocess.run(cmd, check=True)
        print(f"✅ Exported slides to {output_dir}/")
        return output_dir
    except (subprocess.CalledProcessError, FileNotFoundError):
        print("⚠️  ImageMagick not found. Using alternative method...")
        return export_slides_manual(pptx_path, output_dir)

def export_slides_manual(pptx_path, output_dir='slides_export'):
    """
    Manual export method - creates simple representations of slides
    For actual production, use PowerPoint/Keynote export
    """
    os.makedirs(output_dir, exist_ok=True)

    prs = Presentation(pptx_path)

    print(f"Creating slide representations...")
    print(f"⚠️  Note: For production quality, export slides from PowerPoint/Keynote")
    print(f"   This method creates simple previews only.\n")

    for i, slide in enumerate(prs.slides):
        # Create blank image with dark background
        img = Image.new('RGB', (1920, 1080), color=(43, 51, 61))
        draw = ImageDraw.Draw(img)

        # Try to extract and render text (simplified)
        y_offset = 100
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                try:
                    # Use system font
                    font = ImageFont.truetype("/System/Library/Fonts/Helvetica.ttc", 40)
                except:
                    font = ImageFont.load_default()

                # Wrap text
                text = shape.text[:100]  # Limit length
                draw.text((100, y_offset), text, fill=(229, 229, 229), font=font)
                y_offset += 80

        # Save image
        output_path = f'{output_dir}/slide_{i:02d}.png'
        img.save(output_path)
        print(f"  [{i+1:2d}/15] Created: {output_path}")

    print(f"\n✅ Created slide representations in {output_dir}/")
    return output_dir

def create_video_from_images(image_dir, output_video='vibe-manager-video.mp4',
                             with_audio=False, audio_file=None):
    """
    Create video from exported slide images using FFmpeg
    """
    if not check_ffmpeg():
        print("\n❌ Error: FFmpeg is not installed")
        print("\nTo install FFmpeg:")
        print("  macOS:   brew install ffmpeg")
        print("  Linux:   sudo apt-get install ffmpeg")
        print("  Windows: Download from https://ffmpeg.org/download.html")
        return None

    print(f"\n{'='*60}")
    print(f"Creating Video with FFmpeg")
    print(f"{'='*60}\n")

    config = VIDEO_CONFIG

    # FFmpeg command for creating video from images
    # Each image shown for slide_duration seconds with crossfade transition

    cmd = [
        'ffmpeg',
        '-y',  # Overwrite output file
        '-framerate', f'1/{config["slide_duration"]}',  # One frame per slide_duration
        '-pattern_type', 'glob',
        '-i', f'{image_dir}/slide_*.png',
        '-vf', f'scale={config["width"]}:{config["height"]}:force_original_aspect_ratio=decrease,'
               f'pad={config["width"]}:{config["height"]}:(ow-iw)/2:(oh-ih)/2,'
               f'format=yuv420p',
        '-c:v', config['codec'],
        '-b:v', config['bitrate'],
        '-r', str(config['fps']),
        '-pix_fmt', 'yuv420p',
    ]

    # Add audio if provided
    if with_audio and audio_file and os.path.exists(audio_file):
        cmd.extend(['-i', audio_file, '-c:a', 'aac', '-b:a', '192k'])
        print(f"Adding background audio: {audio_file}")

    cmd.append(output_video)

    print(f"FFmpeg command: {' '.join(cmd)}\n")
    print(f"Generating video...")
    print(f"  Resolution: {config['width']}x{config['height']}")
    print(f"  Duration per slide: {config['slide_duration']}s")
    print(f"  FPS: {config['fps']}")
    print(f"  Codec: {config['codec']}")
    print(f"  Bitrate: {config['bitrate']}\n")

    try:
        result = subprocess.run(cmd, capture_output=True, text=True)

        if result.returncode == 0:
            # Get file size
            size_mb = os.path.getsize(output_video) / (1024 * 1024)
            print(f"\n{'='*60}")
            print(f"✅ Video created successfully!")
            print(f"{'='*60}")
            print(f"  Output: {output_video}")
            print(f"  Size: {size_mb:.2f} MB")

            # Get video duration
            duration_cmd = [
                'ffprobe',
                '-v', 'error',
                '-show_entries', 'format=duration',
                '-of', 'default=noprint_wrappers=1:nokey=1',
                output_video
            ]
            duration_result = subprocess.run(duration_cmd, capture_output=True, text=True)
            if duration_result.returncode == 0:
                duration = float(duration_result.stdout.strip())
                print(f"  Duration: {duration:.1f}s ({duration/60:.1f} minutes)")

            print(f"{'='*60}\n")
            return output_video
        else:
            print(f"\n❌ Error creating video:")
            print(result.stderr)
            return None

    except Exception as e:
        print(f"\n❌ Error: {str(e)}")
        return None

def create_video_with_transitions(image_dir, output_video='vibe-manager-video-transitions.mp4'):
    """
    Create video with smooth crossfade transitions between slides
    """
    if not check_ffmpeg():
        print("\n❌ Error: FFmpeg is not installed")
        return None

    print(f"\n{'='*60}")
    print(f"Creating Video with Crossfade Transitions")
    print(f"{'='*60}\n")

    config = VIDEO_CONFIG

    # Get list of images
    images = sorted([f for f in os.listdir(image_dir) if f.endswith('.png')])
    num_slides = len(images)

    if num_slides == 0:
        print("❌ No images found in directory")
        return None

    print(f"Found {num_slides} slides")

    # Build complex FFmpeg filter for crossfade transitions
    # Each slide: slide_duration - transition_duration visible, then transition_duration fade
    slide_time = config['slide_duration']
    trans_time = config['transition_duration']
    hold_time = slide_time - trans_time

    # Create filter_complex for crossfade
    filter_parts = []
    inputs = []

    for i, img in enumerate(images):
        inputs.extend(['-loop', '1', '-t', str(slide_time), '-i', f'{image_dir}/{img}'])

    # Build crossfade chain
    if num_slides == 1:
        filter_complex = '[0:v]'
    else:
        filter_complex = '[0:v][1:v]xfade=transition=fade:duration={}:offset={}'.format(
            trans_time, hold_time
        )

        for i in range(2, num_slides):
            prev_offset = hold_time + (slide_time * (i - 1))
            filter_complex = f'[{filter_complex}][{i}:v]xfade=transition=fade:duration={trans_time}:offset={prev_offset}'

    # Final scaling and formatting
    filter_complex = f'{filter_complex}[v];[v]scale={config["width"]}:{config["height"]}:' \
                    f'force_original_aspect_ratio=decrease,' \
                    f'pad={config["width"]}:{config["height"]}:(ow-iw)/2:(oh-ih)/2,' \
                    f'format=yuv420p[out]'

    cmd = ['ffmpeg', '-y'] + inputs + [
        '-filter_complex', filter_complex,
        '-map', '[out]',
        '-c:v', config['codec'],
        '-b:v', config['bitrate'],
        '-r', str(config['fps']),
        '-pix_fmt', 'yuv420p',
        output_video
    ]

    print(f"Generating video with crossfade transitions...")
    print(f"  Slide duration: {slide_time}s")
    print(f"  Transition duration: {trans_time}s")
    print(f"  Total duration: ~{(slide_time * num_slides) - (trans_time * (num_slides - 1)):.1f}s\n")

    try:
        result = subprocess.run(cmd, capture_output=True, text=True)

        if result.returncode == 0:
            size_mb = os.path.getsize(output_video) / (1024 * 1024)
            print(f"\n{'='*60}")
            print(f"✅ Video with transitions created!")
            print(f"{'='*60}")
            print(f"  Output: {output_video}")
            print(f"  Size: {size_mb:.2f} MB")
            print(f"{'='*60}\n")
            return output_video
        else:
            print(f"\n❌ Error creating video:")
            print(result.stderr)
            return None

    except Exception as e:
        print(f"\n❌ Error: {str(e)}")
        return None

def print_manual_instructions():
    """Print manual instructions for creating video"""
    print(f"\n{'='*60}")
    print(f"Manual Video Creation Instructions")
    print(f"{'='*60}\n")

    print("METHOD 1: PowerPoint (Windows/Mac)")
    print("-" * 40)
    print("1. Open vibe-manager-professional.pptx in PowerPoint")
    print("2. Go to: File > Export > Create a Video")
    print("3. Settings:")
    print("   - Quality: Full HD (1920x1080)")
    print("   - Seconds spent on each slide: 8")
    print("   - Use recorded timings: No")
    print("4. Click 'Create Video' and save as MP4\n")

    print("METHOD 2: Keynote (Mac)")
    print("-" * 40)
    print("1. Open presentation in Keynote (converts from PPTX)")
    print("2. Go to: File > Export To > Movie...")
    print("3. Settings:")
    print("   - Resolution: 1080p")
    print("   - Format: H.264")
    print("   - Slide Duration: 8 seconds")
    print("   - Transitions: Include")
    print("4. Click 'Next' and save\n")

    print("METHOD 3: Google Slides (Web)")
    print("-" * 40)
    print("1. Upload PPTX to Google Drive")
    print("2. Open with Google Slides")
    print("3. Go to: File > Download > Microsoft PowerPoint (.pptx)")
    print("4. Then use PowerPoint or other tool to create video\n")

    print("METHOD 4: LibreOffice Impress (Free, All Platforms)")
    print("-" * 40)
    print("1. Install: https://www.libreoffice.org/")
    print("2. Open presentation in Impress")
    print("3. Go to: Slide Show > Slide Show Settings")
    print("4. Export to PDF, then convert PDF to video with tools\n")

    print(f"{'='*60}\n")

def main():
    """Main execution"""
    import argparse

    parser = argparse.ArgumentParser(description='Generate video from PowerPoint presentation')
    parser.add_argument('--pptx', default='vibe-manager-professional.pptx',
                       help='Input PowerPoint file')
    parser.add_argument('--output', default='vibe-manager-video.mp4',
                       help='Output video file')
    parser.add_argument('--with-transitions', action='store_true',
                       help='Use crossfade transitions (requires more processing)')
    parser.add_argument('--audio', type=str,
                       help='Background audio file (MP3/WAV)')
    parser.add_argument('--slide-duration', type=float, default=8,
                       help='Seconds per slide (default: 8)')
    parser.add_argument('--manual', action='store_true',
                       help='Show manual instructions only')

    args = parser.parse_args()

    if args.manual:
        print_manual_instructions()
        return

    # Update config
    VIDEO_CONFIG['slide_duration'] = args.slide_duration

    # Check dependencies
    print(f"\n{'='*60}")
    print(f"Vibe Manager Video Generator")
    print(f"{'='*60}\n")
    print("Checking dependencies...")

    has_ffmpeg = check_ffmpeg()
    print(f"  FFmpeg: {'✅ Installed' if has_ffmpeg else '❌ Not found'}")

    if not has_ffmpeg:
        print("\n⚠️  FFmpeg is required for video generation")
        print_manual_instructions()
        return

    # Step 1: Export slides as images
    export_dir = 'slides_export'
    export_slides_as_images(args.pptx, export_dir)

    # Check if we have images
    images = [f for f in os.listdir(export_dir) if f.endswith('.png')]
    if not images:
        print("\n⚠️  No slides exported. Showing manual instructions:")
        print_manual_instructions()
        return

    # Step 2: Create video
    if args.with_transitions:
        create_video_with_transitions(export_dir, args.output)
    else:
        create_video_from_images(export_dir, args.output,
                                with_audio=bool(args.audio),
                                audio_file=args.audio)

    print("\n💡 Tips for better results:")
    print("  - Export slides from PowerPoint/Keynote for highest quality")
    print("  - Add voiceover narration for engagement")
    print("  - Use background music (royalty-free)")
    print("  - Consider adding animated text overlays")
    print("  - Upload to YouTube/Vimeo for hosting\n")

if __name__ == '__main__':
    main()
